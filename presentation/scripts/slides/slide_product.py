"""Slide 02 — Product: Dräger SmartCoffee"""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from common import *

from pptx.util import Inches

# Absolute path to the product image
_IMG = os.path.join(os.path.dirname(__file__), "draeger-smart-coffee-machine.png")


def add_slide(prs):
    slide = new_slide(prs)
    hdr(slide, "Dräger SmartCoffee")

    # Centered product image — takes up the bulk of the slide
    img_w = 6.5
    img_h = 5.0
    img_l = (SW - img_w) / 2
    img_t = 1.05
    slide.shapes.add_picture(_IMG, Inches(img_l), Inches(img_t), Inches(img_w), Inches(img_h))

    # Slogan below the image
    txt(slide, "Coffee saves lives.",
        0, img_t + img_h + 0.15, SW, 0.75,
        sz=30, bold=True, col=TEXT, align=PP_ALIGN.CENTER)

    notes(slide,
          "Meet the Dräger SmartCoffee — our demo domain, because coffee saves lives. "
          "A coffee machine in a hospital operating room that writes an audit trail of all its operations to ensure patient safety. "
          "This is a requirement by compliance regulations, but also a great example of an event-sourced system. ")
    return slide
