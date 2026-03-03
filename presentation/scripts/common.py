"""
common.py — Shared theme constants and helper functions for all slide modules.
"""

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme ─────────────────────────────────────────────────────────────────────
BG      = "1a1a2e"   # slide background
PANEL   = "16213e"   # content panel / card fill
BORDER  = "0f3460"   # card border / code background
ACCENT  = "c8a96e"   # coffee gold — headings, highlights
TEXT    = "e0e0e0"   # body text
MUTED   = "888888"   # secondary / metadata text
SUCCESS = "4caf50"   # green
DANGER  = "ef9a9a"   # red
FONT    = "Segoe UI"
MONO    = "Consolas"

SW = 13.33  # slide width in inches (16:9)
SH = 7.5    # slide height in inches


# ── Colour helper ─────────────────────────────────────────────────────────────
def rgb(hex_str: str) -> RGBColor:
    return RGBColor(
        int(hex_str[0:2], 16),
        int(hex_str[2:4], 16),
        int(hex_str[4:6], 16),
    )


# ── Slide factory ─────────────────────────────────────────────────────────────
def new_slide(prs):
    """Add a blank slide with the dark background colour."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    return slide


# ── Shape helpers ─────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    """Add a rectangle (auto-shape type 1). Dimensions in inches."""
    sp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = rgb(line)
        sp.line.width = lw
    else:
        sp.line.fill.background()
    return sp


def txt(slide, text, l, t, w, h,
        fn=None, sz=18, bold=False, col=TEXT,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a text box. Dimensions in inches."""
    fn = fn or FONT
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = fn
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(col)
    return tb


def code(slide, text, l, t, w, h, sz=11):
    """Add a dark code block with monospace text."""
    rect(slide, l, t, w, h, fill=BORDER, line=ACCENT)
    tb = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 0.12),
        Inches(w - 0.3), Inches(h - 0.24),
    )
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = MONO
        r.font.size = Pt(sz)
        r.font.color.rgb = rgb(ACCENT)
    return tb


def notes(slide, text: str):
    """Set speaker notes on a slide."""
    slide.notes_slide.notes_text_frame.text = text


# ── Layout helpers ────────────────────────────────────────────────────────────
def hdr(slide, text, sz=28):
    """Render a standard gold-on-dark header band at the top of the slide."""
    rect(slide, 0, 0, SW, 0.85, fill=PANEL)
    txt(slide, text, 0.3, 0.08, SW - 0.6, 0.72,
        sz=sz, bold=True, col=ACCENT)


def footer(slide, text, col=MUTED, sz=13):
    """Render a thin accent line + italic footer text at the bottom."""
    rect(slide, 0, 6.8, SW, 0.05, fill=ACCENT)
    txt(slide, text, 0.3, 6.88, SW - 0.6, 0.55,
        sz=sz, col=col, align=PP_ALIGN.CENTER, italic=True)


def card(slide, l, t, w, h, title, body, tsz=15, bsz=12):
    """Render a PANEL-filled card with a bold ACCENT title and body text."""
    rect(slide, l, t, w, h, fill=PANEL, line=BORDER)
    txt(slide, title, l + 0.1, t + 0.1, w - 0.2, 0.42,
        sz=tsz, bold=True, col=ACCENT)
    txt(slide, body, l + 0.1, t + 0.57, w - 0.2, h - 0.67,
        sz=bsz, col=TEXT, wrap=True)
